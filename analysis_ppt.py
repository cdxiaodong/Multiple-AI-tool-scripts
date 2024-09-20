from pptx import Presentation
import openai
import argparse

# 初始系统消息
system_message = """
你是一名信息安全专家 帮我分析我提供的ppt中的内容并分别生成标题和 摘要
"""

# 读取PPT文件并提取文本内容
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    slides_text = []

    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + " "
        slides_text.append(slide_text.strip())

    return slides_text

# 使用OpenAI API生成标题和摘要
def generate_title_and_summary_with_openai(text, api_key):
    openai.api_key = api_key
    messages = [
        {"role": "system", "content": system_message},
        {"role": "user", "content": text}
    ]

    response = openai.ChatCompletion.create(
        model="gpt-4-o",
        messages=messages
    )
    result = response.choices[0].message.content.strip()
    return result

# 生成Markdown内容
def generate_markdown(slides_text, api_key):
    markdown_content = ""

    for i, text in enumerate(slides_text):
        # 生成标题和摘要
        result = generate_title_and_summary_with_openai(text, api_key)
        lines = result.split("\n")
        if lines:
            title = lines[0].strip()
            summary = "\n".join(lines[1:]).strip()

            # 生成Markdown内容
            markdown_content += f"## {title}\n\n"
            markdown_content += f"{summary}\n\n"

    return markdown_content

# 保存Markdown文件
def save_markdown_file(markdown_content, output_path):
    with open(output_path, "w", encoding="utf-8") as md_file:
        md_file.write(markdown_content)

# 主函数
def main(ppt_path, output_path, api_key):
    slides_text = extract_text_from_ppt(ppt_path)
    markdown_content = generate_markdown(slides_text, api_key)
    save_markdown_file(markdown_content, output_path)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Analyze PPT content and generate summaries using OpenAI API.")
    parser.add_argument("ppt_path", type=str, help="Path to the PPT file")
    parser.add_argument("output_path", type=str, help="Path to the output Markdown file")
    parser.add_argument("api_key", type=str, help="OpenAI API key")

    args = parser.parse_args()

    main(args.ppt_path, args.output_path, args.api_key)
